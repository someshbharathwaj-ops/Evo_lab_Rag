"""
Document loaders for the Evo Lab RAG ingestion pipeline.

Supported formats:
    - PDF  (.pdf)  — via PyMuPDF (fitz)         [required: PyMuPDF>=1.24]
    - PPTX (.pptx) — via python-pptx             [optional: python-pptx>=0.6]
    - DOCX (.docx) — via python-docx             [optional: python-docx>=1.1]
    - Image (.png, .jpg, .jpeg, .bmp, .tiff, .webp)
                   — via pytesseract + Pillow     [optional: pytesseract, Pillow]

Each loader returns the same standard format:
    List[Dict] where each dict has:
        - "text"     (str)  : extracted/cleaned text content
        - "metadata" (dict) : at minimum {"source": str, "page": int | None}

Only `load_pdf` requires an installed package (PyMuPDF). All other loaders
raise a clear ImportError if their optional dependency is missing.

Usage:
    from ingestion.loaders.loaders import load_document

    docs = load_document("slides.pptx")   # auto-dispatches by extension
    docs = load_pdf("report.pdf")
    docs = load_pptx("slides.pptx")
    docs = load_docx("notes.docx")
    docs = load_image("diagram.png")
"""

from __future__ import annotations

from pathlib import Path
from typing import Dict, List


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _clean_text(text: str) -> str:
    """Normalise whitespace and strip control characters."""
    text = text.replace("\r\n", " ").replace("\n", " ").replace("\r", " ").replace("\t", " ")
    # Strip non-printable control chars (keep normal ASCII + Unicode letters)
    text = "".join(ch for ch in text if ch >= " " or ch in "\t\n\r")
    return " ".join(text.split()).strip()


def _make_doc(text: str, source: str, page: int | None) -> Dict:
    return {
        "text": text,
        "metadata": {
            "source": source,
            "page": page,
        },
    }


# ---------------------------------------------------------------------------
# PDF loader  (PyMuPDF — required dependency, already in requirements.txt)
# ---------------------------------------------------------------------------

def load_pdf(file_path: str) -> List[Dict]:
    """
    Load a PDF and extract per-page text using PyMuPDF.

    Returns:
        List of dicts with 'text' and 'metadata' (source, page).
    """
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise ImportError(
            "PyMuPDF is required for PDF loading. Install it with: pip install PyMuPDF>=1.24"
        ) from exc

    source = str(Path(file_path).resolve())
    documents: List[Dict] = []
    try:
        with fitz.open(file_path) as doc:
            for page_number in range(doc.page_count):
                page = doc.load_page(page_number)
                cleaned = _clean_text(page.get_text())
                if cleaned:
                    documents.append(_make_doc(cleaned, source, page_number + 1))
    except Exception as exc:
        raise RuntimeError(f"Failed to load PDF '{file_path}': {exc}") from exc
    return documents


# ---------------------------------------------------------------------------
# PPTX loader  (python-pptx — optional)
# ---------------------------------------------------------------------------

def load_pptx(file_path: str) -> List[Dict]:
    """
    Load a PowerPoint (.pptx) file and extract text slide-by-slide.

    Each slide becomes one document entry.  Text is gathered from all
    shapes that contain a text frame (titles, text boxes, tables, etc.).

    Returns:
        List of dicts with 'text' and 'metadata' (source, page=slide_number).
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Pt  # noqa: F401 — validates the install
    except ImportError as exc:
        raise ImportError(
            "python-pptx is required for PPTX loading. Install it with: pip install python-pptx>=0.6"
        ) from exc

    source = str(Path(file_path).resolve())
    documents: List[Dict] = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                # Plain text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text)
                        if line.strip():
                            parts.append(line.strip())
                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            combined = _clean_text(" ".join(parts))
            if combined:
                documents.append(_make_doc(combined, source, slide_num))
    except Exception as exc:
        raise RuntimeError(f"Failed to load PPTX '{file_path}': {exc}") from exc
    return documents


# ---------------------------------------------------------------------------
# DOCX loader  (python-docx — optional)
# ---------------------------------------------------------------------------

def load_docx(file_path: str) -> List[Dict]:
    """
    Load a Word (.docx) file and extract text paragraph-by-paragraph,
    grouping content into logical page-equivalent chunks of ~50 paragraphs.

    Because DOCX files have no hard page boundaries, the entire document is
    returned as a single entry (page=1) unless it is very long, in which
    case it is split into chunks of 50 paragraphs each.

    Returns:
        List of dicts with 'text' and 'metadata' (source, page).
    """
    try:
        from docx import Document  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-docx is required for DOCX loading. Install it with: pip install python-docx>=1.1"
        ) from exc

    source = str(Path(file_path).resolve())
    documents: List[Dict] = []
    PARAGRAPHS_PER_PAGE = 50
    try:
        doc = Document(file_path)
        paragraphs = [_clean_text(p.text) for p in doc.paragraphs if p.text.strip()]

        # Also extract table cell text
        for table in doc.tables:
            for row in table.rows:
                cells = [_clean_text(cell.text) for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))

        if not paragraphs:
            return documents

        # Group into page-like chunks
        for page_num, start in enumerate(range(0, len(paragraphs), PARAGRAPHS_PER_PAGE), start=1):
            chunk_text = " ".join(paragraphs[start : start + PARAGRAPHS_PER_PAGE])
            if chunk_text:
                documents.append(_make_doc(chunk_text, source, page_num))
    except Exception as exc:
        raise RuntimeError(f"Failed to load DOCX '{file_path}': {exc}") from exc
    return documents


# ---------------------------------------------------------------------------
# Image loader  (pytesseract + Pillow — optional, requires Tesseract OCR)
# ---------------------------------------------------------------------------

SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}


def load_image(file_path: str) -> List[Dict]:
    """
    Extract text from an image file using Tesseract OCR via pytesseract.

    Requires both the pytesseract Python package AND the Tesseract binary to
    be installed on the system:
        pip install pytesseract Pillow
        # Windows: download from https://github.com/UB-Mannheim/tesseract/wiki
        # Linux:   sudo apt install tesseract-ocr

    Returns:
        List with a single dict: 'text' (OCR result) and 'metadata' (source, page=1).
    """
    try:
        import pytesseract  # type: ignore[import]
        from PIL import Image  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "pytesseract and Pillow are required for image loading. "
            "Install with: pip install pytesseract Pillow\n"
            "Also install the Tesseract binary for your OS."
        ) from exc

    source = str(Path(file_path).resolve())
    try:
        img = Image.open(file_path)
        raw_text = pytesseract.image_to_string(img, lang="eng")
        cleaned = _clean_text(raw_text)
        if not cleaned:
            return []
        return [_make_doc(cleaned, source, 1)]
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from image '{file_path}': {exc}") from exc


# ---------------------------------------------------------------------------
# Universal dispatcher — load_document()
# ---------------------------------------------------------------------------

_EXTENSION_MAP = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".bmp": load_image,
    ".tiff": load_image,
    ".tif": load_image,
    ".webp": load_image,
}


def load_document(file_path: str) -> List[Dict]:
    """
    Auto-dispatch to the correct loader based on file extension.

    Supported extensions: .pdf, .pptx, .docx, .png, .jpg, .jpeg,
                          .bmp, .tiff, .tif, .webp

    Raises:
        ValueError: if the file extension is not supported.
        FileNotFoundError: if the file does not exist.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    loader = _EXTENSION_MAP.get(ext)
    if loader is None:
        supported = ", ".join(sorted(_EXTENSION_MAP.keys()))
        raise ValueError(
            f"Unsupported file type '{ext}' for '{file_path}'. "
            f"Supported: {supported}"
        )
    return loader(file_path)
